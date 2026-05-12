from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)   # dark navy – backgrounds, headers
TEAL       = RGBColor(0x00, 0x8B, 0x8B)   # teal – accent boxes
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF9)   # slide background
MID_GREY   = RGBColor(0x6C, 0x75, 0x7D)   # body text
AMBER      = RGBColor(0xFF, 0xA5, 0x00)   # highlight
GREEN      = RGBColor(0x28, 0xA7, 0x45)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color=LIGHT_GREY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def bullet_box(slide, lines, l, t, w, h,
               size=14, color=NAVY, leading_color=TEAL, bg_color=None):
    """Render a list of (bullet_char, text) tuples as a textbox."""
    if bg_color:
        box(slide, l, t, w, h, bg_color)
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (symbol, line) in lines:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        # coloured symbol
        r1 = p.add_run()
        r1.text = symbol + "  "
        r1.font.size  = Pt(size)
        r1.font.bold  = True
        r1.font.color.rgb = leading_color
        # body text
        r2 = p.add_run()
        r2.text = line
        r2.font.size  = Pt(size)
        r2.font.color.rgb = color
        p.space_after = Pt(4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – Title
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)

# Left accent bar
box(s1, 0, 0, 0.18, 7.5, TEAL)

# Title
txt(s1, "QE Story Reviewer", 0.45, 1.6, 9.0, 1.2,
    size=44, bold=True, color=WHITE)

# Subtitle
txt(s1, "AI-powered Quality Engineering in Jira Rovo",
    0.45, 2.9, 9.0, 0.7, size=22, color=RGBColor(0xB0,0xC4,0xDE))

# Divider
box(s1, 0.45, 3.75, 5.5, 0.05, TEAL)

# Three tags
tags = ["Shift-Left QE", "Three Amigos", "BDD Gherkin"]
for i, tag in enumerate(tags):
    bx = box(s1, 0.45 + i * 2.8, 4.1, 2.4, 0.55, TEAL)
    txt(s1, tag, 0.55 + i * 2.8, 4.15, 2.2, 0.45,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Role line
txt(s1, "Testing Specialist  ·  Salesforce Sales Cloud  ·  Agile Team",
    0.45, 5.1, 10.0, 0.5, size=13, color=RGBColor(0x80,0x9A,0xBF), italic=True)

# Right illustration text
txt(s1, "@QE Story Reviewer\nrun three amigos on this story",
    9.8, 3.2, 3.3, 1.2, size=12, color=RGBColor(0x40,0x60,0x80),
    italic=True, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – How It Works (flow)
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)

# Header band
box(s2, 0, 0, 13.33, 1.1, NAVY)
txt(s2, "How It Works", 0.3, 0.15, 8.0, 0.8,
    size=28, bold=True, color=WHITE)
txt(s2, "Paste a Jira ticket  →  get QE analysis + Gherkin in seconds",
    0.3, 0.55, 10.0, 0.45, size=13, color=RGBColor(0xB0,0xC4,0xDE))

# ── Flow steps ───────────────────────────────────────────────────────────────
steps = [
    ("1", "Open Ticket\nin Jira",           NAVY),
    ("2", "@QE Story\nReviewer",            TEAL),
    ("3", "Agent reads\nSkill + Knowledge", NAVY),
    ("4", "QE Analysis\nposted to Jira",    GREEN),
]

step_w, step_h = 2.2, 1.4
gap = 0.55
start_x = 0.5
y = 2.0

for i, (num, label, col) in enumerate(steps):
    x = start_x + i * (step_w + gap)
    box(s2, x, y, step_w, step_h, col)
    txt(s2, num, x + 0.12, y + 0.08, 0.4, 0.4,
        size=20, bold=True, color=AMBER)
    txt(s2, label, x + 0.12, y + 0.45, step_w - 0.2, 0.85,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = x + step_w + 0.1
        txt(s2, "→", ax, y + 0.45, 0.35, 0.5,
            size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ── What triggers it ─────────────────────────────────────────────────────────
txt(s2, "What triggers the agent", 0.5, 3.75, 5.5, 0.45,
    size=15, bold=True, color=NAVY)
triggers = [
    ("▸", "Paste a user story, epic, or bug ticket"),
    ("▸", '"review this story" / "is this ready for dev?"'),
    ("▸", '"run three amigos on this story"'),
    ("▸", '"generate Gherkin for this AC"'),
    ("▸", '"analyse this defect"'),
]
bullet_box(s2, triggers, 0.5, 4.2, 5.8, 2.6, size=13, color=MID_GREY, leading_color=TEAL)

# ── What you type ─────────────────────────────────────────────────────────────
box(s2, 6.8, 3.75, 6.2, 3.3, RGBColor(0xE8, 0xF4, 0xF4))
txt(s2, "Example in Rovo chat", 7.0, 3.85, 5.8, 0.4,
    size=13, bold=True, color=TEAL)
txt(s2,
    "@QE Story Reviewer\n\n"
    "Title: Export transaction history as CSV\n\n"
    "As a customer I want to export my transactions\n"
    "so I can import into accounting software.\n\n"
    "AC1: Given I am logged in, when I click\n"
    "Export CSV, then a file downloads.",
    7.0, 4.3, 5.8, 2.5,
    size=11, color=NAVY, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – Three Amigos
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)

box(s3, 0, 0, 13.33, 1.1, NAVY)
txt(s3, "Three Amigos — Simulated by the Agent", 0.3, 0.15, 10.0, 0.8,
    size=28, bold=True, color=WHITE)
txt(s3, "One ticket in  ·  Three perspectives out  ·  Agreed Gherkin + DoD",
    0.3, 0.55, 10.0, 0.45, size=13, color=RGBColor(0xB0,0xC4,0xDE))

cols = [
    {
        "title": "Business / PO",
        "icon":  "🏢",
        "color": RGBColor(0x1B,0x4F,0x72),
        "items": [
            "User value clear & measurable?",
            "All business rules written down?",
            "Compliance / regulatory rules?",
            "Salesforce licence tier / record type?",
            "Approval process in scope?",
            "→ 3–7 questions for the PO",
        ]
    },
    {
        "title": "Salesforce Developer",
        "icon":  "⚙️",
        "color": RGBColor(0x0E,0x6B,0x6B),
        "items": [
            "Apex vs Flow — right tool chosen?",
            "Governor limit risks identified?",
            "Deployment sequence defined?",
            "External API contracts agreed?",
            "Test class coverage plan?",
            "→ 3–7 questions for Dev",
        ]
    },
    {
        "title": "Tester / QA",
        "icon":  "✅",
        "color": RGBColor(0x14,0x5A,0x32),
        "items": [
            "INVEST check (6 dimensions)",
            "AC rewritten as Gherkin",
            "Edge cases & negative paths",
            "NFRs: perf, security, a11y",
            "Risk register (L × I)",
            "→ Gherkin feature file",
        ]
    },
]

col_w  = 3.7
col_h  = 4.6
start_x = 0.45
y_top  = 1.35

for i, col in enumerate(cols):
    x = start_x + i * (col_w + 0.45)
    # header band
    box(s3, x, y_top, col_w, 0.65, col["color"])
    txt(s3, col["title"], x + 0.12, y_top + 0.1, col_w - 0.2, 0.5,
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # body
    box(s3, x, y_top + 0.65, col_w, col_h - 0.65, RGBColor(0xF0,0xF5,0xF5))
    items = [("•", it) for it in col["items"]]
    bullet_box(s3, items, x + 0.15, y_top + 0.8,
               col_w - 0.25, col_h - 0.9,
               size=12, color=NAVY, leading_color=col["color"])

# Synthesis arrow + box
box(s3, 0.45, 6.1, 12.43, 0.05, TEAL)
txt(s3, "▼  Synthesis: Gaps table  ·  Refined Gherkin (annotated with business rule + layer)  ·  Agreed DoD  ·  Verdict",
    0.45, 6.2, 12.4, 0.55,
    size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – What Your Team Gets
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)

box(s4, 0, 0, 13.33, 1.1, NAVY)
txt(s4, "What Your Team Gets", 0.3, 0.15, 9.0, 0.8,
    size=28, bold=True, color=WHITE)
txt(s4, "Every output is Jira-pasteable — paste as a comment on the ticket",
    0.3, 0.55, 10.0, 0.45, size=13, color=RGBColor(0xB0,0xC4,0xDE))

# ── Left: output list ─────────────────────────────────────────────────────────
outputs = [
    ("✅", "Verdict: Ready for Dev / Needs Refinement / Blocked"),
    ("📋", "INVEST check across 6 dimensions"),
    ("⚠️", "Risk register — Likelihood × Impact"),
    ("🔍", "Missing scenarios: edge cases, negative paths, NFRs"),
    ("🧪", "Full BDD Gherkin feature file with layer + priority tags"),
    ("❓", "3–7 open questions for PO / Dev"),
    ("🏁", "Tailored Definition of Done checklist"),
    ("👥", "Three Amigos: Business + Dev + QA perspectives merged"),
]
txt(s4, "Output sections", 0.45, 1.25, 6.0, 0.4,
    size=14, bold=True, color=NAVY)
bullet_box(s4, outputs, 0.45, 1.7, 6.2, 5.4,
           size=13, color=NAVY, leading_color=TEAL)

# ── Right: Gherkin sample ─────────────────────────────────────────────────────
box(s4, 7.1, 1.25, 5.9, 5.85, RGBColor(0x1E,0x1E,0x2E))
txt(s4, "Sample Gherkin output", 7.25, 1.35, 5.6, 0.4,
    size=12, bold=True, color=TEAL)

gherkin = (
    "Feature: Export Transaction History\n\n"
    "  Background:\n"
    "    Given I am logged in as a Customer\n\n"
    "  # Layer: e2e | Priority: P0\n"
    "  Scenario: Download CSV for filtered range\n"
    "    Given I have filtered transactions\n"
    "      by date range Jan–Mar 2026\n"
    "    When I click Export CSV\n"
    "    Then a file downloads named\n"
    "      transactions_2026-03-31.csv\n"
    "    And it contains date, amount,\n"
    "      currency, description, balance\n\n"
    "  # Layer: api | Priority: P1\n"
    "  Scenario: Export blocked over 100K rows\n"
    "    Given I have 150,000 transactions\n"
    "    When I click Export CSV\n"
    "    Then I see: 'Export limited to 100K\n"
    "      rows. Refine your date filter.'"
)
txt(s4, gherkin, 7.25, 1.8, 5.6, 5.0,
    size=10, color=RGBColor(0xCC,0xFF,0xCC), italic=True)

# ── Bottom banner ─────────────────────────────────────────────────────────────
box(s4, 0, 7.1, 13.33, 0.4, NAVY)
txt(s4,
    "Rovo Agent  ·  Knowledge: 7 workflow & reference files  ·  "
    "Trigger: paste any Jira ticket or type @QE Story Reviewer",
    0.3, 7.13, 12.7, 0.32,
    size=10, color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
out = r"c:\Users\honra\Claude\QE-Framework\QE_Story_Reviewer_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
